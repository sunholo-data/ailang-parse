#!/usr/bin/env python3
"""Verification loop for generated Office documents.

Checks that generated files are valid at multiple levels:
  Level 1: Structural (ZIP + XML well-formedness)
  Level 2: Library (python-docx, python-pptx, openpyxl open without errors)
  Level 4: Roundtrip (generate → parse through DocParse → verify blocks)

Usage:
  uv run --with python-pptx --with openpyxl --with python-docx benchmarks/verify_generated.py
  uv run --with python-pptx --with openpyxl --with python-docx benchmarks/verify_generated.py --fix  # regenerate after fixes
"""

import json
import os
import posixpath
import subprocess
import sys
import tempfile
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path

REPO_DIR = Path(__file__).parent.parent
EXAMPLES_DIR = REPO_DIR / "data" / "examples"

# ── Level 1: Structural Checks ──────────────────────────────────────────────

def verify_structure(path: Path) -> tuple[list[str], list[str]]:
    """Check ZIP structure and XML well-formedness. Returns (errors, warnings)."""
    errors = []
    warnings = []
    if not zipfile.is_zipfile(path):
        return ([f"Not a valid ZIP file"], [])

    with zipfile.ZipFile(path) as zf:
        names = zf.namelist()

        # Check required entries by format
        ext = path.suffix.lower()
        if ext == ".docx":
            required = ["[Content_Types].xml", "_rels/.rels", "word/document.xml"]
        elif ext == ".pptx":
            required = ["[Content_Types].xml", "_rels/.rels", "ppt/presentation.xml"]
        elif ext == ".xlsx":
            required = ["[Content_Types].xml", "_rels/.rels", "xl/workbook.xml"]
        elif ext in (".odt", ".odp", ".ods"):
            required = ["mimetype", "META-INF/manifest.xml", "content.xml"]
        else:
            required = []

        for req in required:
            if req not in names:
                errors.append(f"Missing required entry: {req}")

        # Check XML well-formedness for all .xml and .rels entries
        for name in names:
            if name.endswith(".xml") or name.endswith(".rels"):
                try:
                    data = zf.read(name).decode("utf-8")
                    ET.fromstring(data)
                except ET.ParseError as e:
                    errors.append(f"XML parse error in {name}: {e}")
                except Exception as e:
                    errors.append(f"Error reading {name}: {e}")

        # PPTX-specific: check theme exists
        if ext == ".pptx":
            has_theme = any("theme" in n for n in names)
            if not has_theme:
                errors.append("Missing ppt/theme/theme1.xml (Keynote requires this)")

        # ODF-specific: check mimetype is first entry and uncompressed
        # NOTE: AILANG's createArchive always compresses. This is a known limitation
        # that needs a stdlib change (store flag). Downgraded to warning.
        if ext in (".odt", ".odp", ".ods"):
            if names and names[0] != "mimetype":
                warnings.append(f"mimetype should be first ZIP entry, got: {names[0]}")
            info = zf.getinfo("mimetype") if "mimetype" in names else None
            if info and info.compress_type != zipfile.ZIP_STORED:
                warnings.append("mimetype entry is compressed (ODF spec wants uncompressed) — AILANG stdlib limitation")

    return (errors, warnings)


# ── Level 2: Library Validation ──────────────────────────────────────────────

def verify_library(path: Path) -> list[str]:
    """Open with Python libraries and check for errors/warnings."""
    errors = []
    ext = path.suffix.lower()
    import warnings

    try:
        if ext == ".docx":
            from docx import Document
            doc = Document(str(path))
            if len(doc.paragraphs) == 0:
                errors.append("DOCX has 0 paragraphs")
            errors.extend(_docx_table_grid_errors(doc))
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            if len(prs.slides) == 0:
                errors.append("PPTX has 0 slides")
            for i, slide in enumerate(prs.slides):
                if len(slide.shapes) == 0:
                    errors.append(f"Slide {i+1} has 0 shapes")
        elif ext == ".xlsx":
            from openpyxl import load_workbook
            with warnings.catch_warnings(record=True) as w:
                warnings.simplefilter("always")
                wb = load_workbook(str(path))
                for warning in w:
                    errors.append(f"Warning: {warning.message}")
                if len(wb.sheetnames) == 0:
                    errors.append("XLSX has 0 sheets")
    except Exception as e:
        errors.append(f"{type(e).__name__}: {e}")

    return errors



def _docx_table_grid_errors(doc) -> list[str]:
    """Every row must span exactly the columns w:tblGrid declares.

    ECMA-376 defines the grid as part of the table, and a row that does not
    add up to it is malformed however forgiving the reader: python-docx raises
    on cell access and Word offers to repair. This caught a real defect —
    horizontal merge continuations were emitted as extra <w:vMerge> cells on
    top of the gridSpan that already covered them, so one row of a 4-column
    table was 7 grid units wide. Opening the file cleanly did not catch it,
    because LibreOffice tolerates it.
    """
    from docx.oxml.ns import qn

    errors = []
    for ti, table in enumerate(doc.tables):
        grid = table._tbl.find(qn("w:tblGrid"))
        if grid is None:
            errors.append(f"table {ti}: no w:tblGrid")
            continue
        declared = len(grid.findall(qn("w:gridCol")))
        for ri, row in enumerate(table._tbl.findall(qn("w:tr"))):
            width = 0
            for tc in row.findall(qn("w:tc")):
                pr = tc.find(qn("w:tcPr"))
                span = 1
                if pr is not None:
                    gs = pr.find(qn("w:gridSpan"))
                    if gs is not None:
                        span = int(gs.get(qn("w:val")))
                width += span
            if width != declared:
                errors.append(
                    f"table {ti} row {ri}: spans {width} grid columns, "
                    f"tblGrid declares {declared}")
        # Cell iteration is what actually fails on a malformed grid.
        try:
            for row in table.rows:
                _ = [c.text for c in row.cells]
        except Exception as e:
            errors.append(f"table {ti}: cell iteration failed — {type(e).__name__}: {e}")
    return errors


# ── Level 2b: DOCX package wiring ────────────────────────────────────────────

def verify_docx_parts(path: Path) -> list[str]:
    """Assert every declared part is reachable through the relationship graph.

    OPC resolves parts via relationships, so a part that is written and declared
    in [Content_Types].xml but never related to is invisible to Word. That is
    exactly how word/styles.xml shipped for months while every heading rendered
    as body text: the XML was correct, the package wiring was not.

    Checked generically rather than per-feature so that numbering.xml,
    comments.xml and header/footer parts are covered by the same assertion the
    day they are added.
    """
    errors = []
    with zipfile.ZipFile(path) as z:
        names = set(z.namelist())

        declared = set()
        root = ET.fromstring(z.read("[Content_Types].xml"))
        for ov in root.findall("{http://schemas.openxmlformats.org/package/2006/content-types}Override"):
            declared.add(ov.get("PartName").lstrip("/"))

        # Resolve every relationship target to a package-absolute part name.
        #
        # Every _rels part in the package is walked, and targets are resolved
        # against the directory that owns them — including "../" hops. Checking
        # only _rels/.rels and word/_rels/document.xml.rels was enough while
        # every package here was generated from scratch; a reference doc brings
        # customXml/_rels/item1.xml.rels and Target="../customXml/item1.xml"
        # with it, and both read as broken wiring under the old resolution.
        reachable = {"word/document.xml"}
        rel_ns = "{http://schemas.openxmlformats.org/package/2006/relationships}Relationship"
        for rels_name in sorted(n for n in names if n.endswith(".rels")):
            owner_dir = posixpath.dirname(posixpath.dirname(rels_name))
            for rel in ET.fromstring(z.read(rels_name)).findall(rel_ns):
                if rel.get("TargetMode") == "External":
                    continue
                target = rel.get("Target") or ""
                if target.startswith("/"):
                    resolved = target.lstrip("/")
                else:
                    resolved = posixpath.normpath(posixpath.join(owner_dir, target))
                reachable.add(resolved)

        for part in sorted(declared - reachable):
            errors.append(f"part declared but unreachable (no relationship): {part}")
        for part in sorted(p for p in reachable if p not in names and not p.startswith("docProps/app")):
            errors.append(f"relationship points at missing part: {part}")

        # The user-visible symptom of the above: headings falling back to body
        # text. If the body asks for a Heading style, a reader must resolve it.
        doc_xml = z.read("word/document.xml").decode("utf-8", "replace")
        if 'w:val="Heading' in doc_xml:
            try:
                from docx import Document
                styles = {p.style.name for p in Document(str(path)).paragraphs if p.style}
                if not any(s.startswith("Heading") for s in styles):
                    errors.append("document uses Heading styles but none resolve (orphaned styles.xml?)")
            except Exception as e:
                errors.append(f"heading-style check failed: {type(e).__name__}: {e}")

    return errors


# ── Level 4: Roundtrip Verification ──────────────────────────────────────────

def verify_roundtrip(path: Path) -> list[str]:
    """Parse the generated file through DocParse and check output."""
    errors = []
    ext = path.suffix.lower()

    # Skip formats DocParse can't re-parse (ODP, ODS have basic parsers)
    if ext in (".html", ".htm"):
        return []  # HTML roundtrip already tested

    with tempfile.TemporaryDirectory() as tmpdir:
        try:
            result = subprocess.run(
                ["ailang", "run", "--entry", "main", "--caps", "IO,FS,Env",
                 "--max-recursion-depth", "50000",
                 "docparse/main.ail", str(path)],
                capture_output=True, text=True, cwd=str(REPO_DIR),
                timeout=30,
                env={**os.environ, "DOCPARSE_OUTPUT_DIR": tmpdir},
            )
            if result.returncode != 0:
                stderr_lines = result.stderr.strip().split("\n")[-3:]
                errors.append(f"DocParse failed (exit {result.returncode}): {' '.join(stderr_lines)}")
                return errors

            # Check output mentions blocks
            output = result.stdout
            if "Blocks:   0" in output:
                errors.append("Roundtrip produced 0 blocks")

            # Extract block count from output
            for line in output.split("\n"):
                if line.strip().startswith("Blocks:"):
                    count = line.strip().split()[-1]
                    if count == "0":
                        errors.append("Roundtrip produced 0 blocks")

        except subprocess.TimeoutExpired:
            errors.append("Roundtrip timed out (30s)")
        except Exception as e:
            errors.append(f"Roundtrip error: {e}")

    return errors


# ── Main ─────────────────────────────────────────────────────────────────────

# Sources deliberately chosen to carry SectionBlocks (DOCX comments/text boxes/
# section breaks, XLSX sheet sections, HTML header/footer) — the shape that made
# every non-presentation source hang on ->PPTX and ->ODP. Markdown has none,
# which is why the advertised "notes.md --convert slides.pptx" path stayed green
# while six other conversions were dead.
CONVERSION_MATRIX_SOURCES = [
    "inline_formatting.docx",
    "unstructured_test.xlsx",
    "test.html",
    "test.md",
]
CONVERSION_TARGETS = ["html", "docx", "pptx", "xlsx", "odt", "odp", "ods", "md", "qmd"]


# ── Reference-doc (Quarto-style template) verification ──────────────────────

REFERENCE_DOC_SOURCE = """---
title: Reference Doc Check
author: DocParse
---

# Top Heading

## Second Level

Body paragraph with a [link](https://example.com/x).

- bullet one
- bullet two

1. ordered one
2. ordered two

| Col A | Col B |
|---|---|
| a1 | b1 |
| a2 | b2 |
"""

# Both templates are in-repo and chosen for what they stress:
#   docx-hdrftr.docx  even/default headers AND footers on the body sectPr, an
#                     attribute-bearing <w:sectPr w:rsidR=...>, no numbering.xml,
#                     and a styles.xml missing Heading1/2 and ListParagraph
#   comments.docx     comments.xml + commentsExtended.xml + people.xml, which the
#                     merge has to drop together with their rels and overrides
REFERENCE_DOC_TEMPLATES = ["docx-hdrftr.docx", "comments.docx"]

# Parts the merge is expected to rewrite. Anything else in the template must
# come out byte-for-byte, because "the template's look survived" IS the feature.
_REF_REGENERATED = {
    "[Content_Types].xml", "_rels/.rels", "docProps/core.xml",
    "word/document.xml", "word/_rels/document.xml.rels",
    "word/styles.xml", "word/numbering.xml",
}
_REF_DROPPED = {
    "word/comments.xml", "word/commentsExtended.xml",
    "word/commentsIds.xml", "word/people.xml",
}


def _body_sect_pr(doc_xml: str) -> str:
    """The template's own body-level <w:sectPr>, revision spans removed."""
    import re
    inner = doc_xml.rsplit("</w:body>", 1)[0]
    inner = re.sub(r"<w:sectPrChange\b.*?</w:sectPrChange>", "", inner, flags=re.S)
    m = list(re.finditer(r"<w:sectPr[ >].*?</w:sectPr>", inner, re.S))
    return m[-1].group(0) if m else ""


def _reference_doc_errors(template: Path, out: Path) -> list[str]:
    import re
    errors = []
    with zipfile.ZipFile(template) as tz, zipfile.ZipFile(out) as oz:
        tnames, onames = set(tz.namelist()), set(oz.namelist())

        # 1. Nothing from the template is silently lost, and what is carried is
        #    carried unchanged — this is what keeps the fonts and theme intact.
        for n in sorted(tnames - _REF_REGENERATED - _REF_DROPPED):
            if n.endswith("/"):
                continue
            if n not in onames:
                errors.append(f"template part dropped: {n}")
            elif tz.read(n) != oz.read(n):
                errors.append(f"carried template part was modified: {n}")

        # 2. The page setup, headers and footers come from the template. Losing
        #    this is the failure that still opens cleanly and looks like ours.
        t_sect = _body_sect_pr(tz.read("word/document.xml").decode("utf-8", "replace"))
        o_sect = _body_sect_pr(oz.read("word/document.xml").decode("utf-8", "replace"))
        if not t_sect:
            errors.append(f"fixture {template.name} has no body sectPr to compare")
        elif o_sect != t_sect:
            errors.append(f"body sectPr not lifted from template\n"
                          f"       template: {t_sect[:160]}\n"
                          f"       output:   {o_sect[:160]}")

        # 3. The template's own comment parts leave together with their wiring.
        for n in sorted(_REF_DROPPED & tnames):
            if n in onames:
                errors.append(f"template comment part should have been dropped: {n}")

        # 4. Our list numbering must not land on a numId the template defines,
        #    which renders as the template's list rather than a bullet.
        if "word/numbering.xml" in tnames:
            tpl_nums = set(re.findall(r'<w:num w:numId="(\d+)"', tz.read("word/numbering.xml").decode()))
            used = set(re.findall(r'<w:numId w:val="(\d+)"', oz.read("word/document.xml").decode()))
            clash = tpl_nums & used
            if clash:
                errors.append(f"body uses numId(s) the template already defines: {sorted(clash)}")

        # 5. Our media must not overwrite the template's.
        tpl_media = {n for n in tnames if n.startswith("word/media/")}
        for n in tpl_media:
            if n in onames and tz.read(n) != oz.read(n):
                errors.append(f"template media overwritten: {n}")

        # 6. The template's styles survive and ours only fill gaps.
        tpl_ids = set(re.findall(r'w:styleId="([^"]+)"', tz.read("word/styles.xml").decode()))
        out_ids = set(re.findall(r'w:styleId="([^"]+)"', oz.read("word/styles.xml").decode()))
        missing = tpl_ids - out_ids
        if missing:
            errors.append(f"template styleIds lost: {sorted(missing)[:6]}")
        for needed in ("ListParagraph", "Hyperlink"):
            if needed not in out_ids:
                errors.append(f"styleId referenced by the body is undefined: {needed}")

        # 8. Tables bind to the template's table style when it defines a usable
        #    one (anything but the implicit Normal Table every table has).
        #    Hardcoded borders would override the style and defeat the binding,
        #    so their presence under a bound style is itself the defect.
        tpl_styles = tz.read("word/styles.xml").decode("utf-8", "replace")
        tpl_table_ids = set(re.findall(
            r'<w:style w:type="table"[^>]*w:styleId="([^"]+)"', tpl_styles))
        tpl_table_named = dict(re.findall(
            r'<w:style w:type="table"[^>]*w:styleId="([^"]+)"><w:name w:val="([^"]+)"', tpl_styles))
        usable = {sid for sid, name in tpl_table_named.items()
                  if name not in ("Normal Table", "Table Normal")} | \
                 {sid for sid in tpl_table_ids if tpl_table_named.get(sid) not in ("Normal Table", "Table Normal")}
        if usable:
            body = oz.read("word/document.xml").decode("utf-8", "replace")
            style_refs = set(re.findall(r'<w:tblStyle w:val="([^"]+)"', body))
            if not style_refs:
                errors.append("template defines table style(s) but body binds to none")
            else:
                dangling = style_refs - tpl_table_ids
                if dangling:
                    errors.append(f"tblStyle references undefined styleId(s): {sorted(dangling)}")
            if "<w:tblBorders>" in body:
                errors.append("hardcoded tblBorders present under a bound table style "
                              "(direct formatting overrides the style)")

    # 7. The generic package wiring assertions apply here too — a dangling
    #    relationship or an Override naming an absent part is what makes Word
    #    refuse the file outright.
    errors.extend(verify_docx_parts(out))

    # 8. And it has to open.
    try:
        from docx import Document
        d = Document(str(out))
        errors.extend(_docx_table_grid_errors(d))
        if not any((p.style.name or "").startswith("Heading") for p in d.paragraphs if p.style):
            errors.append("no heading resolved through the template's styles")
    except Exception as e:
        errors.append(f"python-docx could not open the output: {type(e).__name__}: {e}")

    return errors


def verify_reference_doc() -> bool:
    """--reference-doc must apply a template's look without breaking the package.

    Read BACK, not just opened: every defect this guards against — a lost
    letterhead, a bullet rendering as the template's numbered list, a clobbered
    logo — produces a file that opens perfectly and is wrong.
    """
    repo = Path(__file__).resolve().parent.parent
    test_dir = repo / "data" / "test_files"
    print("── reference doc ──")
    ok = True
    with tempfile.TemporaryDirectory() as tmp:
        src = Path(tmp) / "refdoc_source.md"
        src.write_text(REFERENCE_DOC_SOURCE, encoding="utf-8")
        for name in REFERENCE_DOC_TEMPLATES:
            template = test_dir / name
            if not template.exists():
                print(f"  L6 RefDoc:     SKIP ({name} missing)")
                continue
            out = Path(tmp) / f"styled__{template.stem}.docx"
            r = subprocess.run(
                [str(repo / "bin" / "docparse"), str(src), "--convert", str(out),
                 "--reference-doc", str(template)],
                capture_output=True, text=True, errors="replace",
                cwd=str(repo), timeout=120,
            )
            if r.returncode != 0 or not out.exists() or out.stat().st_size == 0:
                lines = ((r.stdout or "") + (r.stderr or "")).splitlines()
                err = next((l for l in reversed(lines) if "rror" in l), lines[-1] if lines else "no output")
                print(f"  L6 RefDoc:     FAIL ({name}: {err.strip()[:120]})")
                ok = False
                continue
            errs = _reference_doc_errors(template, out)
            if errs:
                print(f"  L6 RefDoc:     FAIL ({name})")
                for e in errs:
                    print(f"     ⚠ {e}")
                ok = False
            else:
                print(f"  L6 RefDoc:     PASS ({name})")

        # A reference doc that cannot be read must write nothing rather than
        # quietly producing an unstyled file that looks like a success.
        bad_out = Path(tmp) / "should_not_exist.docx"
        r = subprocess.run(
            [str(repo / "bin" / "docparse"), str(src), "--convert", str(bad_out),
             "--reference-doc", str(test_dir / "ailang_formats.csv")],
            capture_output=True, text=True, errors="replace", cwd=str(repo), timeout=120,
        )
        if bad_out.exists():
            print("  L6 RefDoc:     FAIL (non-DOCX reference still produced output)")
            ok = False
        else:
            print("  L6 RefDoc:     PASS (non-DOCX reference refused)")

        # --reference-section: a multi-section template has one sectPr per
        # section plus the body one, and the body one is the LAST. The in-repo
        # templates are single-section, so the stage synthesizes a two-section
        # variant of docx-hdrftr.docx whose FIRST section (paragraph-level,
        # distinct margins and header ref) must be liftable by number.
        base = test_dir / "docx-hdrftr.docx"
        if base.exists():
            two = _synthesize_two_section_template(base, Path(tmp) / "two-section.docx")
            if two is None:
                print("  L6 RefSection: FAIL (could not synthesize two-section template)")
                ok = False
            else:
                ok = _check_reference_section(repo, src, two, Path(tmp)) and ok
    return ok


def _synthesize_two_section_template(base: Path, dest: Path) -> Path | None:
    """docx-hdrftr.docx with a paragraph-level sectPr injected before the body
    one: section 1 = injected (pgMar top=1000, header rId11), section 2 = the
    original body sectPr. Returns the path, or None on any surprise."""
    import re
    try:
        with zipfile.ZipFile(base) as zin:
            doc = zin.read("word/document.xml").decode("utf-8")
            body_sect = re.search(r"<w:sectPr[ >].*</w:sectPr>", doc, re.S)
            if not body_sect:
                return None
            mid = ('<w:p><w:pPr><w:sectPr>'
                   '<w:headerReference w:type="default" r:id="rId11"/>'
                   '<w:pgSz w:w="11906" w:h="16838"/>'
                   '<w:pgMar w:top="1000" w:right="1000" w:bottom="1000" w:left="1000"/>'
                   '</w:sectPr></w:pPr></w:p>')
            newdoc = doc.replace(body_sect.group(0), mid + body_sect.group(0))
            with zipfile.ZipFile(dest, "w", zipfile.ZIP_DEFLATED) as zout:
                for n in zin.namelist():
                    zout.writestr(n, newdoc.encode("utf-8") if n == "word/document.xml" else zin.read(n))
        return dest
    except Exception:
        return None


def _check_reference_section(repo: Path, src: Path, template: Path, tmp: Path) -> bool:
    """--reference-section 1 lifts the FIRST sectPr (injected pgMar top=1000);
    the default (no flag) still lifts the LAST; out of range writes nothing."""
    import re
    ok = True

    def _lifted(out: Path) -> str:
        with zipfile.ZipFile(out) as z:
            doc = z.read("word/document.xml").decode("utf-8", "replace")
        m = list(re.finditer(r"<w:sectPr[ >].*</w:sectPr>", doc, re.S))
        return m[-1].group(0) if m else ""

    out1 = tmp / "sect1.docx"
    r = subprocess.run(
        [str(repo / "bin" / "docparse"), str(src), "--convert", str(out1),
         "--reference-doc", str(template), "--reference-section", "1"],
        capture_output=True, text=True, errors="replace", cwd=str(repo), timeout=120)
    if r.returncode != 0 or not out1.exists():
        print("  L6 RefSection: FAIL (--reference-section 1 errored)")
        return False
    s1 = _lifted(out1)
    if 'w:top="1000"' not in s1:
        print("  L6 RefSection: FAIL (section 1 not lifted — injected pgMar absent)")
        ok = False
    else:
        print("  L6 RefSection: PASS (--reference-section 1 lifts the first sectPr)")

    outd = tmp / "sectdefault.docx"
    subprocess.run(
        [str(repo / "bin" / "docparse"), str(src), "--convert", str(outd),
         "--reference-doc", str(template)],
        capture_output=True, text=True, errors="replace", cwd=str(repo), timeout=120)
    if outd.exists() and 'w:top="1000"' not in _lifted(outd):
        print("  L6 RefSection: PASS (no flag still lifts the last sectPr)")
    else:
        print("  L6 RefSection: FAIL (default no longer lifts the last sectPr)")
        ok = False

    out3 = tmp / "sect3.docx"
    subprocess.run(
        [str(repo / "bin" / "docparse"), str(src), "--convert", str(out3),
         "--reference-doc", str(template), "--reference-section", "3"],
        capture_output=True, text=True, errors="replace", cwd=str(repo), timeout=120)
    if out3.exists():
        print("  L6 RefSection: FAIL (out-of-range section still produced output)")
        ok = False
    else:
        print("  L6 RefSection: PASS (out-of-range refused, nothing written)")
    return ok


def verify_conversion_matrix() -> bool:
    """Every source format must convert to every target without hanging."""
    import subprocess, tempfile

    repo = Path(__file__).resolve().parent.parent
    test_dir = repo / "data" / "test_files"
    print("── conversion matrix ──")
    failures = []
    with tempfile.TemporaryDirectory() as tmp:
        for src in CONVERSION_MATRIX_SOURCES:
            src_path = test_dir / src
            if not src_path.exists():
                continue
            for tgt in CONVERSION_TARGETS:
                out = Path(tmp) / f"{src_path.stem}__{tgt}.{tgt}"
                try:
                    r = subprocess.run(
                        [str(repo / "bin" / "docparse"), str(src_path), "--convert", str(out)],
                        capture_output=True, text=True, errors="replace",
                        cwd=str(repo), timeout=120,
                    )
                except subprocess.TimeoutExpired:
                    failures.append(f"{src} -> {tgt}: timed out"); continue
                if r.returncode != 0 or not out.exists() or out.stat().st_size == 0:
                    # Surface the real error, not whatever happened to print last.
                    lines = ((r.stdout or "") + (r.stderr or "")).splitlines()
                    err = next((l for l in reversed(lines) if "Error" in l or "error" in l),
                               lines[-1] if lines else "no output")
                    failures.append(f"{src} -> {tgt}: {err.strip()[:100]}")
    total = len(CONVERSION_MATRIX_SOURCES) * len(CONVERSION_TARGETS)
    if failures:
        print(f"  L5 Conversions: FAIL ({len(failures)}/{total})")
        for f in failures:
            print(f"     ⚠ {f}")
        return False
    print(f"  L5 Conversions: PASS ({total}/{total})")
    return True


# Hostile inputs whose text is attacker-controlled, paired with the substring
# that must NOT appear in generated output. Each is a real breakout, not a
# lookalike: the payload closes the attribute it is interpolated into.
#
# Added after `<img src='x&quot; onerror=&quot;alert(1)'>` in an input HTML
# document produced a live `onerror="alert(1)"` handler in the generated one.
# html_generator escaped 19 of its 20 interpolation sites; the image src was the
# one that did not, so nothing structural caught it. "Does it open" cannot.
# Targets are per-case, and deliberately not uniform.
#
# HTML is the format where document text is unambiguously data: anything from
# the source that lands in markup is a defect. Markdown and Quarto permit inline
# HTML by design, and docparse is a lossless converter — escaping a <script> the
# author really wrote would itself be a bug. So a raw-tag assertion only applies
# to HTML output, while an attribute breakout is checked everywhere it could
# manifest. (QMD front matter is separately safe: the generator backslash-escapes
# quotes, so a title containing " does not break out of the YAML string.)
INJECTION_CASES = [
    (
        "img_src_breakout.html",
        "<html><body><img src='x&quot; onerror=&quot;alert(1)' alt='p'/></body></html>",
        ['onerror="alert(1)"', "onerror='alert(1)'"],
        ("html", "qmd", "md"),
    ),
    (
        "link_href_breakout.html",
        "<html><body><a href='y&quot; onclick=&quot;alert(2)'>t</a></body></html>",
        ['onclick="alert(2)"', "onclick='alert(2)'"],
        ("html", "qmd", "md"),
    ),
    (
        "title_breakout.html",
        "<html><head><title>a&lt;/title&gt;&lt;script&gt;alert(3)&lt;/script&gt;</title></head>"
        "<body><p>b</p></body></html>",
        ["<script>alert(3)</script>", "</title><script>"],
        ("html",),
    ),
]


def verify_no_injection() -> bool:
    """Attacker-controlled document text must not become markup in the output."""
    import subprocess, tempfile

    repo = Path(__file__).resolve().parent.parent
    print("── injection escaping ──")
    failures = []
    checked = 0
    with tempfile.TemporaryDirectory() as tmp:
        for name, source, forbidden, targets in INJECTION_CASES:
            src = Path(tmp) / name
            src.write_text(source)
            for tgt in targets:
                checked += 1
                out = Path(tmp) / f"{src.stem}__{tgt}.{tgt}"
                r = subprocess.run(
                    [str(repo / "bin" / "docparse"), str(src), "--convert", str(out)],
                    capture_output=True, text=True, errors="replace",
                    cwd=str(repo), timeout=120,
                )
                if r.returncode != 0 or not out.exists():
                    failures.append(f"{name} -> {tgt}: conversion failed")
                    continue
                text = out.read_text(errors="replace")
                for bad in forbidden:
                    if bad in text:
                        failures.append(f"{name} -> {tgt}: emitted {bad!r}")

    if failures:
        print(f"  L6 Injection:  FAIL ({len(failures)})")
        for f in failures:
            print(f"     ⚠ {f}")
        return False
    print(f"  L6 Injection:  PASS ({checked}/{checked} — no attribute or tag breakout)")
    return True


def main():
    print("=== DocParse Generated File Verification ===\n")

    files = sorted(EXAMPLES_DIR.glob("*"))
    files = [f for f in files if f.suffix.lower() in (".docx", ".pptx", ".xlsx", ".odt", ".odp", ".ods", ".html")]

    if not files:
        print(f"No files found in {EXAMPLES_DIR}")
        sys.exit(1)

    all_pass = verify_conversion_matrix()
    print()
    all_pass = verify_no_injection() and all_pass
    print()
    all_pass = verify_reference_doc() and all_pass
    print()
    for path in files:
        print(f"── {path.name} ──")

        # Level 1: Structure
        if path.suffix.lower() in (".html", ".htm"):
            # HTML: just check it's valid XML (XHTML)
            try:
                ET.parse(str(path))
                print(f"  L1 Structure:  PASS (valid XHTML)")
            except ET.ParseError as e:
                print(f"  L1 Structure:  WARN (not valid XHTML: {e})")
        else:
            struct_errors, struct_warnings = verify_structure(path)
            if struct_errors:
                print(f"  L1 Structure:  FAIL")
                for e in struct_errors:
                    print(f"     ⚠ {e}")
                all_pass = False
            elif struct_warnings:
                print(f"  L1 Structure:  WARN")
                for w in struct_warnings:
                    print(f"     ⚠ {w}")
            else:
                print(f"  L1 Structure:  PASS")

        # Level 2: Library
        if path.suffix.lower() in (".docx", ".pptx", ".xlsx"):
            lib_errors = verify_library(path)
            if lib_errors:
                print(f"  L2 Library:    WARN")
                for e in lib_errors:
                    print(f"     ⚠ {e}")
            else:
                print(f"  L2 Library:    PASS")

        # Level 2b: DOCX package wiring (orphaned parts are a FAIL, not a warn —
        # they produce a file that opens fine and silently drops formatting)
        if path.suffix.lower() == ".docx":
            part_errors = verify_docx_parts(path)
            if part_errors:
                print(f"  L2b Parts:     FAIL")
                for e in part_errors:
                    print(f"     ⚠ {e}")
                all_pass = False
            else:
                print(f"  L2b Parts:     PASS")

        # Level 4: Roundtrip (skip HTML — already tested separately)
        if path.suffix.lower() not in (".html", ".htm"):
            rt_errors = verify_roundtrip(path)
            if rt_errors:
                print(f"  L4 Roundtrip:  FAIL")
                for e in rt_errors:
                    print(f"     ⚠ {e}")
                all_pass = False
            else:
                print(f"  L4 Roundtrip:  PASS")

        print()

    if all_pass:
        print("=== ALL CHECKS PASSED ===")
    else:
        print("=== SOME CHECKS FAILED (see above) ===")
        sys.exit(1)


if __name__ == "__main__":
    main()
